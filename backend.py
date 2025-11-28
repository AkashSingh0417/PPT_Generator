import os
import json
import requests
from dotenv import load_dotenv
import google.generativeai as genai
from pptx import Presentation
from layout_style import theme_map
from pptx.dml.color import RGBColor
from pptx.util import Inches
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.util import Pt

load_dotenv()
genai.configure(api_key=os.getenv("GEMINI_API_KEY"))
UNSPLASH_ACCESS_KEY = os.getenv("UNSPLASH_ACCESS_KEY")

#to get image from the unsplash
def get_unsplash_image(query):
    if not UNSPLASH_ACCESS_KEY:
        print("❌ Unsplash API key is not set. Cannot search for images.")
        return None
    try:
        url = "https://api.unsplash.com/search/photos"
        params = {
            "query": query,
            "client_id": UNSPLASH_ACCESS_KEY,
            "orientation": "landscape",
            "per_page": 1
        }
        response = requests.get(url, params=params)
        response.raise_for_status()
        data = response.json()
        if data["results"]:
            image_url = data["results"][0]["urls"]["regular"]
            image_response = requests.get(image_url)
            image_response.raise_for_status()

            temp_image_path = "temp_image.jpg"
            with open(temp_image_path, "wb") as f:
                f.write(image_response.content)
            return temp_image_path
    except Exception as e:
        print(f"❌ Error getting image for query '{query}': {e}")
        return None

# Ask Gemini to create slides
def generate_slides(prompt):
    model = genai.GenerativeModel("gemini-2.5-flash")
    response = model.generate_content(
        f"""Create a 6-slide presentation in hybrid JSON+Markdown format.
        The output must include:
        - A "layout_style" field that MUST always be present. The value MUST be one of: ["Minimalist", "Corporate", "Creative", "Educational", "Technology", "Elegant", "Energetic"]. Pick the one that best fits the topic. Do not invent new ones.
        - The number of bullet points per slide should be fully dynamic 
          (decide based on the topic and importance of that slide).
        - Bullet points should not be too short: 
          each should be an informative sentence (not just 1-2 words).
        - If needed, include 1-2 longer bullets that explain key details or examples.
        - A "slides" list containing slide titles and content.

        Input: {prompt}

        Format:
        {{
          "layout_style": "string",
          "slides": [
            {{
              "title": "## Title",
              "content": "- Bullet point 1\\n- Bullet point 2",
              "image_query": "string" 
            }}
          ]
        }}
        """

    )
    return response.text

# Convert JSON response into PPT

def create_ppt_from_json(json_str, filename, template_path=None):
    data = json.loads(json_str)

    # Load template if provided, else start with blank
    prs = Presentation(template_path) if template_path else Presentation()

    MAX_POINTS = 6
    for slide_data in data["slides"]:
        clean_content = slide_data["content"].replace("-", "").replace("*", "").strip()
        points = [p.strip() for p in clean_content.split(". ") if p.strip()]

        for i in range(0, len(points), MAX_POINTS):
            chunk = points[i:i + MAX_POINTS]

            slide_layout = prs.slide_layouts[1]  # Title + Content from template
            slide = prs.slides.add_slide(slide_layout)

            # Title
            title = slide.shapes.title
            title.text = slide_data["title"].replace("## ", "").strip()

            # Content
            content = slide.placeholders[1]
            text_frame = content.text_frame
            text_frame.clear()
            for point in chunk:
                p = text_frame.add_paragraph()
                p.text = point

    prs.save(filename)


if __name__ == "__main__":
    user_prompt = input("Enter the topic for your PowerPoint presentation:")
    hybrid_output = generate_slides(user_prompt)

    try:
        json_part = hybrid_output[hybrid_output.index("{"):hybrid_output.rindex("}")+1]
        create_ppt_from_json(json_part, "try_presentation.pptx")
    except Exception as e:
        print("⚠️ Error parsing JSON:", e)
        print("Model Output:", hybrid_output)
